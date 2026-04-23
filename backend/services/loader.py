from __future__ import annotations

"""
Document text extraction for RAG.

Why `rglob` under a base path: uploads are organized in nested topic folders; scanning
from a root picks up every supported file without hard-coding layout.

Why per-format loaders: students upload PDFs, slides, and Word docs; one entrypoint
(`load_documents`) keeps RAG code agnostic of file type.

Why lazy imports inside ppt/doc helpers: keeps startup lighter if those deps are missing
until someone actually indexes those types.
"""

import logging
from pathlib import Path
from typing import List, Sequence, Tuple

import fitz

logger = logging.getLogger(__name__)

_TEXT_BLOCK = 0  # PyMuPDF: text vs image blocks—we skip image blocks to avoid OCR noise.
_LINE_Y_TOL = 6.0  # PDF layout often fragments one line into boxes; merging needs loose vertical tolerance.
_MIN_BLOCK_CHARS = 3  # Tiny boxes are usually figure labels, not prose worth embedding.
_MIN_LINE_CHARS = 2  # After merge, single-letter lines are almost always noise, not definitions.


def _normalize_ws(s: str) -> str:
    return " ".join(s.split())


def _blocks_to_page_text(
    blocks: Sequence[Tuple],
    *,
    line_y_tol: float = _LINE_Y_TOL,
    min_block_chars: int = _MIN_BLOCK_CHARS,
    min_line_chars: int = _MIN_LINE_CHARS,
) -> str:
    """
    Turn PyMuPDF block tuples into reading-order lines: merge boxes on the same row,
    drop tiny boxes that are usually slide/diagram noise.
    """
    items: List[Tuple[float, float, str]] = []
    for b in blocks:
        if len(b) < 7:
            continue
        x0, y0, x1, y1, raw, _block_no, btype = b[0], b[1], b[2], b[3], b[4], b[5], b[6]
        if btype != _TEXT_BLOCK:
            continue
        t = _normalize_ws((raw or "").replace("\n", " "))
        if len(t) < min_block_chars:
            continue
        cy = (float(y0) + float(y1)) / 2.0
        items.append((cy, float(x0), t))

    if not items:
        return ""

    items.sort(key=lambda r: (r[0], r[1]))

    out_lines: List[str] = []
    row: List[Tuple[float, str]] = []
    anchor_cy: float | None = None

    def flush_row() -> None:
        nonlocal row, anchor_cy
        if not row:
            return
        row.sort(key=lambda x: x[0])
        line = " ".join(seg for _, seg in row)
        line = _normalize_ws(line)
        if len(line) >= min_line_chars:
            out_lines.append(line)
        row = []
        anchor_cy = None

    for cy, x0, t in items:
        if anchor_cy is None or abs(cy - anchor_cy) <= line_y_tol:
            row.append((x0, t))
            if anchor_cy is None:
                anchor_cy = cy
        else:
            flush_row()
            row.append((x0, t))
            anchor_cy = cy
    flush_row()

    return "\n".join(out_lines)


def load_txt(base: Path) -> List[str]:
    texts: List[str] = []
    for path in base.rglob("*.txt"):
        try:
            raw = path.read_text(encoding="utf-8", errors="replace").strip()
            if raw:
                texts.append(raw)
        except Exception as e:
            logger.warning("Could not read %s: %s", path.name, e)
    return texts


def load_pdfs(base: Path) -> List[str]:
    texts: List[str] = []
    # De-hyphenate wrapped words; blocks mode preserves layout boxes for merging.
    text_flags = fitz.TEXT_DEHYPHENATE

    for path in base.rglob("*.pdf"):
        try:
            with fitz.open(str(path)) as doc:
                for page in doc:
                    blocks = page.get_text("blocks", flags=text_flags)
                    page_text = _blocks_to_page_text(blocks)
                    if page_text.strip():
                        texts.append(page_text.strip())

        except Exception as e:
            logger.warning("Could not read %s: %s", path.name, e)

    logger.debug("Loaded %s PDF page(s) from %s", len(texts), base)
    return texts


def load_ppt(base: Path) -> List[str]:
    texts: List[str] = []
    for pattern in ("*.pptx", "*.ppt"):
        for path in base.rglob(pattern):
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                for slide in prs.slides:
                    slide_text = " ".join(
                        shape.text for shape in slide.shapes if hasattr(shape, "text")
                    ).strip()
                    if slide_text:
                        texts.append(slide_text)
            except Exception as e:
                logger.warning("Could not read %s: %s", path.name, e)
    return texts


def load_doc(base: Path) -> List[str]:
    texts: List[str] = []
    for pattern in ("*.docx", "*.doc"):
        for path in base.rglob(pattern):
            try:
                from docx import Document
                doc = Document(str(path))
                for para in doc.paragraphs:
                    txt = para.text.strip()
                    if txt:
                        texts.append(txt)
            except Exception as e:
                logger.warning("Could not read %s: %s", path.name, e)
    return texts


def load_documents(uploads_dir: str) -> List[str]:
    base = Path(uploads_dir)
    if not base.exists():
        return []
    logger.debug("Loading documents from %s", uploads_dir)
    documents: List[str] = []
    documents.extend(load_pdfs(base))
    documents.extend(load_txt(base))
    documents.extend(load_ppt(base))
    documents.extend(load_doc(base))
    return documents

